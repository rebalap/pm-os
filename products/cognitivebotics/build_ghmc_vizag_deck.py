"""
Build the GHMC Vizag pilot executive summary deck (10 slides) using the
Cognitivebotics branded PPTX library. Embeds the CBT logo on cover,
content slides, and the closing slide.
"""

import sys, os

# Resolve tools path — works whether run from the Mac (Read tool path) or the
# Linux sandbox (where the workspace is mounted under /sessions/.../mnt/pm-os).
_TOOLS_CANDIDATES = [
    '/Users/prahladrebala/Documents/pm-os/tools',
    '/sessions/epic-loving-tesla/mnt/pm-os/tools',
]
for _p in _TOOLS_CANDIDATES:
    if os.path.isdir(_p):
        sys.path.insert(0, _p)
        break

from cb_pptx import CB_Pptx, W, H, M_LEFT
from pptx.util import Inches

_OUTPUTS_CANDIDATES = [
    '/Users/prahladrebala/Library/Application Support/Claude/local-agent-mode-sessions/d9897f46-568a-47a4-a3a7-eed438edad95/6047fde2-bfb3-48de-a252-11e38becc248/local_438ac185-8aa9-47eb-a163-669bdc891ca7/outputs',
    '/sessions/epic-loving-tesla/mnt/outputs',
]
_OUTPUTS = next(p for p in _OUTPUTS_CANDIDATES if os.path.isdir(p))
LOGO_WHITE = os.path.join(_OUTPUTS, 'cb_logo_white.png')
LOGO_TEAL  = os.path.join(_OUTPUTS, 'cb_logo_teal.png')

_PRODUCT_CANDIDATES = [
    '/Users/prahladrebala/Documents/pm-os/products/cognitivebotics',
    '/sessions/epic-loving-tesla/mnt/cognitivebotics',
]
_PRODUCT_DIR = next(p for p in _PRODUCT_CANDIDATES if os.path.isdir(p))
OUT_PATH = os.path.join(_PRODUCT_DIR, 'GHMC_Vizag_Pilot_Executive_Summary.pptx')

# Logo aspect ratio from SVG: 187 x 25 → ratio ~7.48:1


def add_logo(slide, variant='white', position='top-right'):
    """Place CBT logo on a slide."""
    path = LOGO_WHITE if variant == 'white' else LOGO_TEAL
    # Logo dimensions tuned for each position
    if position == 'cover':
        # Bigger logo top-left on cover
        slide.shapes.add_picture(
            path,
            Inches(0.8), Inches(0.6),
            width=Inches(2.4),
        )
    elif position == 'top-right':
        # Small logo top-right corner of content slides — sits over teal title bar
        slide.shapes.add_picture(
            path,
            W - Inches(2.0), Inches(0.32),
            width=Inches(1.4),
        )
    elif position == 'closing':
        # Bottom-right of closing slide
        slide.shapes.add_picture(
            path,
            W - Inches(2.6), H - Inches(0.7),
            width=Inches(1.8),
        )


prs = CB_Pptx()

# ── Slide 1: Cover ──────────────────────────────────────────────────────────
prs.cover(
    "Closing the Autism Intervention Gap",
    subtitle="GHMC Vizag Pilot — Executive Summary  |  AI-Powered Gamified Home Learning",
    label="Pilot Implementation Report",
)
add_logo(prs._prs.slides[-1], variant='white', position='cover')

# ── Slide 2: The Problem ────────────────────────────────────────────────────
prs.content(
    "The Problem: An Intervention Access Gap",
    bullets=[
        "**Early intensive behavioral intervention works** — but requires 20–40 hours/week of trained-therapist time",
        "**India faces a severe workforce shortage** — fewer than ~500 BCBAs nationally; RCI-licensed special educators are the primary workforce",
        "**Geographic disparities** restrict access; most families cannot reach a therapy center daily",
        "**The home is the missed multiplier** — between clinic sessions, learning is unstructured and ad-hoc, often guided only by WhatsApp",
        "**Result:** Children fall short of the intervention dose research shows they need",
    ],
)
add_logo(prs._prs.slides[-1], variant='white', position='top-right')

# ── Slide 3: Our Solution ───────────────────────────────────────────────────
prs.content(
    "Our Solution: A Digital Extension of Clinical Therapy",
    bullets=[
        "**AI-powered gamified learning platform** for children with autism, ADHD, and other learning differences (ages 2–18)",
        "**Therapist web app** — assigns Individualized Learning Plans (ILPs) drawn from 12 skill domains",
        "**Child-facing app (mobile / iPad)** — gamified delivery of ABA-based instruction at home",
        "**Automated measurement** — eye gaze, speech, and pose detection score every trial without manual data collection",
        "**Designed for non-BCBA workforces** and parents with minimal clinical training",
    ],
)
add_logo(prs._prs.slides[-1], variant='white', position='top-right')

# ── Slide 4: Scientific Foundation ──────────────────────────────────────────
prs.table_slide(
    "Scientific Foundation",
    headers=["Evidence Area", "Established Finding", "How the Platform Applies It"],
    rows=[
        ["Early Intervention",      "Brain plasticity drives developmental gains when intervention starts early",                 "Targets early learning skills in critical 2–18 age window"],
        ["ABA-Based EIBI",          "Meta-analyses show IQ, language, and adaptive behavior gains",                                "Built on reinforcement, prompting, shaping, mastery progression"],
        ["Accessibility Barriers",  "BCBA / RBT shortage limits intervention dose, especially in LMICs",                          "Extends structured learning beyond the clinic, irrespective of location"],
        ["Parent-Implemented",      "Caregiver-led practice increases dose and supports skill generalization",                    "Supports guided home practice with caregiver-friendly UI"],
        ["Tech-Assisted",           "Digital tools complement therapist-delivered programs",                                       "Delivers structured digital activities as a high-fidelity complement"],
        ["Gamified Instruction",    "Game elements act as reinforcement when tied to instructional contingencies",                "Levels, rewards, and interactive feedback reinforce correct responses"],
        ["AI-Powered Measurement",  "ML enables precision treatment via automated behavioral measurement",                         "Response detection, scoring, eye gaze / speech / pose tracking"],
    ],
    col_widths=[2.5, 4.8, 4.6],
)
add_logo(prs._prs.slides[-1], variant='white', position='top-right')

# ── Slide 5: Pilot Design ───────────────────────────────────────────────────
prs.content(
    "Pilot Design — GHMC Vizag",
    bullets=[
        "**Setting:** Real-world service delivery — Vizag pilot cohort, India",
        "**Participants:** 20 children with autism, evaluated at baseline, midline, and endline",
        "**Assessment:** 55-item parent-report questionnaire administered at each timepoint",
        "**Intervention:** Individualized Learning Plan per child — 2–12 Learning Objectives (LOs), mean = 5.9",
        "**Primary metrics:** In-app LO mastery, parent-reported new-skill acquisition, engagement (days played)",
        "**Feasibility layer:** Parent feasibility survey (N=8) capturing UX, integration, and perceived progress",
    ],
)
add_logo(prs._prs.slides[-1], variant='white', position='top-right')

# ── Slide 6: Headline Results ───────────────────────────────────────────────
prs.table_slide(
    "Headline Results",
    headers=["Metric", "Result", "What It Means"],
    rows=[
        ["Families reporting new skills at endline",      "18 / 20  (90%)",   "Near-universal real-world impact"],
        ["Total new skills reported by parents",          "144",              "Across communication, play, daily living"],
        ["Children with ≥1 LO mastered in-app",           "13 / 20  (65%)",   "Two-thirds reached clinical mastery threshold"],
        ["Total LOs mastered (across cohort)",            "55 of 118",        "46.6% overall mastery rate"],
        ["Highest-gaining child (parent-reported)",       "32 new skills",    "16 days played, 5 LOs mastered"],
        ["Children showing no change",                    "2 / 20",           "Both had only 2 days of engagement"],
    ],
    col_widths=[5.2, 3.0, 3.7],
)
add_logo(prs._prs.slides[-1], variant='white', position='top-right')

# ── Slide 7: Engagement → Outcomes ──────────────────────────────────────────
prs.two_column(
    "Engagement Is the Catalyst — A Clear Dose-Response",
    left_label="High Engagement (≥8 days, n=10)",
    right_label="Low Engagement (<8 days, n=10)",
    left=[
        "**16.9** avg days played",
        "**5.3** avg LOs mastered in-app",
        "**9.5** avg new skills reported by parents",
        "**26×** more in-app mastery than the low-engagement group",
    ],
    right=[
        "**3.6** avg days played",
        "**0.2** avg LOs mastered in-app",
        "**4.9** avg new skills reported by parents",
        "Engagement — not demographics — is the modifiable lever",
    ],
)
add_logo(prs._prs.slides[-1], variant='white', position='top-right')

# Add a footnote with correlation strength on the same slide
from pptx.util import Pt
from cb_pptx import _add_textbox
import cb_brand as B
last = prs._prs.slides[-1]
_add_textbox(
    last,
    M_LEFT, H - Inches(0.85), W - Inches(1.4), Inches(0.4),
    text="Internal validity: days played → LOs mastered  r = 0.92, p < 0.001 (strong positive). "
         "Engagement → parent-reported gains  r = 0.45 (moderate positive).",
    size=10, color=B.GRAY, font=B.FONT_BODY, italic=True, align='left',
)

# ── Slide 8: Real-World Skill Gains ─────────────────────────────────────────
prs.two_column(
    "Real-World Skills Gained — Where Impact Showed Up",
    left_label="Mastery by Developmental Domain (n=55 LOs)",
    right_label="Most-Reported New Skills at Endline",
    left=[
        "**Communication & Object ID — 44%**",
        "**Attention & Visual — 31%**",
        "Other developmental domains — 25%",
        "Top LO across cohort: *Learn to Look* — mastered by 12 of 20 children (60%)",
    ],
    right=[
        "Watches moving objects",
        "Responds when name is called",
        "Points to wanted items",
        "Follows simple instructions",
        "Stacks blocks, rings, or cups",
        "Chooses 1 item when given 2 options",
    ],
)
add_logo(prs._prs.slides[-1], variant='white', position='top-right')

# ── Slide 9: Parent Feasibility ─────────────────────────────────────────────
prs.table_slide(
    "Parent Feasibility — High Acceptance, Low Burden",
    headers=["Dimension", "Mean (1–5)", "Caregiver Signal"],
    rows=[
        ["Home Integration",   "4.44",  "Pacing manageable; no added family stress"],
        ["User Experience",    "4.06",  "Setup and instructions clear"],
        ["Child Engagement",   "4.00",  "Consistent interest, willing to participate"],
        ["Perceived Progress", "4.00",  "Observed skill gains and generalization"],
    ],
    col_widths=[3.5, 2.0, 6.4],
)
# Add qualitative callout under the table
last = prs._prs.slides[-1]
_add_textbox(
    last,
    M_LEFT, Inches(4.6), W - Inches(1.4), Inches(2.4),
    text="Usage pattern: families use the platform 4.2 days/week, in 10–15 min micro-sessions. "
         "50% of parents reported only 'sometimes' needing to support the child — an emerging "
         "level of learner independence.\n\n"
         "Qualitative themes: parents observed gains in eye contact, speech patterns, and behavioral "
         "regulation. Animal-sound and visual-tracking activities were the strongest engagement drivers.",
    size=12, color=B.BLACK, font=B.FONT_BODY, align='left',
)
add_logo(prs._prs.slides[-1], variant='white', position='top-right')

# ── Slide 10: Closing — Conclusions & Next Steps ────────────────────────────
prs.closing(
    "Conclusions & Next Steps",
    bullets=[
        "**90% of families** reported real-world skill gains — strong preliminary clinical signal",
        "**Engagement is the lever** — high-engagement children mastered 26× more LOs",
        "**In-app metrics validate as a clinical proxy** — aligned with parent-reported gains (r ≈ 0.42)",
        "**Feasible at home** — 4.44/5 home integration, no caregiver-burden penalty",
        "**Ready to scale** — APIs in place for HMIS / ABDM integration into public health systems",
    ],
)
add_logo(prs._prs.slides[-1], variant='white', position='closing')

# Save
prs.save(OUT_PATH)
print(f'Slides: {len(prs._prs.slides)}')
