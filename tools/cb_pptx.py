"""
cb_pptx.py — Cognitivebotics branded PowerPoint library (python-pptx).

QUICK START
-----------
    from cb_pptx import CB_Pptx

    prs = CB_Pptx()
    prs.cover("Discovery Readout: Center Lifecycle", subtitle="April 2026")
    prs.section_divider("What We Set Out to Learn")
    prs.content("Key Finding 1", bullets=[
        "Special educators record trial data on paper during live sessions",
        "One-handed constraint makes accurate recording physically difficult",
        "Supervisor sees data 1–2 weeks after sessions",
    ])
    prs.two_column(
        "Implication",
        left=["Paper creates transcription errors", "Delayed program updates"],
        right=["Digital tool must be ≤ 2 taps", "Real-time supervisor visibility"],
    )
    prs.table_slide("Hypothesis Tracker",
        headers=["#", "Hypothesis", "Status"],
        rows=[["H-01", "Paper data sheets in live sessions", "Open"]],
    )
    prs.save("discovery-readout.pptx")

CLI
---
    python3 cb_pptx.py input.md output.pptx
    python3 cb_pptx.py input.md output.pptx --title "Discovery Readout"
"""

from __future__ import annotations
import re, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

_dir = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _dir)
import cb_brand as B


# ── Color helpers ─────────────────────────────────────────────────────────────

def _rgb(h: str) -> RGBColor:
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _fill_solid(shape, hex_color: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_color)

def _no_fill(shape):
    shape.fill.background()

def _no_line(shape):
    shape.line.fill.background()


# ── Text helpers ──────────────────────────────────────────────────────────────

def _run(tf_para, text: str, bold: bool = False, italic: bool = False,
         size: int = None, color: str = B.BLACK, font: str = B.FONT_BODY):
    run = tf_para.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.name   = font
    if size:
        run.font.size = Pt(size)
    run.font.color.rgb = _rgb(color)
    return run

def _set_para_align(para, align: str = 'left'):
    mapping = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
               'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY}
    para.alignment = mapping.get(align, PP_ALIGN.LEFT)

def _add_textbox(slide, left, top, width, height,
                 text: str = '', bold: bool = False, italic: bool = False,
                 size: int = B.SZ_BODY, color: str = B.BLACK,
                 font: str = B.FONT_BODY, align: str = 'left',
                 word_wrap: bool = True) -> object:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    _set_para_align(p, align)
    _run(p, text, bold=bold, italic=italic, size=size, color=color, font=font)
    return txBox

def _add_label_rect(slide, left, top, width, height,
                    text: str, bg_color: str = B.TEAL,
                    text_color: str = B.WHITE, size: int = 24,
                    font: str = B.FONT_HEADING, bold: bool = True,
                    align: str = 'left'):
    """Add a filled rectangle with centered text inside."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height)
    _fill_solid(shape, bg_color)
    _no_line(shape)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_para_align(p, align)
    _run(p, text, bold=bold, size=size, color=text_color, font=font)
    return shape

def _inline_runs(tf_para, text: str, size: int = B.SZ_BODY,
                  color: str = B.BLACK, font: str = B.FONT_BODY):
    """Parse **bold** and *italic* markdown and add formatted runs."""
    pattern = r'(\*\*[^*]+\*\*|\*[^*\n]+\*|[^*]+)'
    for m in re.finditer(pattern, text):
        chunk = m.group(0)
        if chunk.startswith('**') and chunk.endswith('**') and len(chunk) > 4:
            _run(tf_para, chunk[2:-2], bold=True, size=size, color=color, font=font)
        elif chunk.startswith('*') and chunk.endswith('*') and len(chunk) > 2:
            _run(tf_para, chunk[1:-1], italic=True, size=size, color=color, font=font)
        else:
            _run(tf_para, chunk, size=size, color=color, font=font)


# ── Slide dimensions ──────────────────────────────────────────────────────────

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

M_LEFT  = Inches(0.7)
M_TOP   = Inches(0.5)
M_RIGHT = Inches(0.7)
CONTENT_W = W - M_LEFT - M_RIGHT
TITLE_BAR_H = Inches(0.9)
CONTENT_TOP = M_TOP + TITLE_BAR_H + Inches(0.1)
CONTENT_H   = H - CONTENT_TOP - Inches(0.5)


# ── CB_Pptx class ─────────────────────────────────────────────────────────────

class CB_Pptx:
    """
    Cognitivebotics branded PowerPoint builder.

    Call slide methods in the order you want them to appear, then `.save(path)`.
    """

    def __init__(self):
        self._prs = Presentation()
        self._prs.slide_width  = W
        self._prs.slide_height = H
        # Use blank layout (index 6 in default template) for full control
        self._blank = self._prs.slide_layouts[6]

    def _new_slide(self):
        slide = self._prs.slides.add_slide(self._blank)
        # Remove any placeholder shapes from the blank layout
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        return slide

    def _title_bar(self, slide, title: str, bg: str = B.TEAL,
                   text_color: str = B.WHITE) -> None:
        """Draw teal title bar at top of slide."""
        bar = slide.shapes.add_shape(1, 0, 0, W, TITLE_BAR_H)
        _fill_solid(bar, bg)
        _no_line(bar)
        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        from pptx.util import Pt as _Pt
        p.space_before = _Pt(8)
        _run(p, title, bold=True, size=24, color=text_color,
             font=B.FONT_HEADING)
        # Left padding via indent
        tf.margin_left = Inches(0.25)
        tf.margin_top  = Inches(0.18)

    def _footer(self, slide, text: str = B.PRODUCT_NAME) -> None:
        """Subtle footer at bottom of slide."""
        _add_textbox(slide,
                     M_LEFT, H - Inches(0.35), CONTENT_W, Inches(0.3),
                     text=text, size=8, color=B.GRAY,
                     font=B.FONT_BODY, align='right')

    # ── Slide types ───────────────────────────────────────────────────────────

    def cover(self, title: str, subtitle: str = '',
              label: str = '') -> 'CB_Pptx':
        """
        Full cover slide: gradient band + title + subtitle.
        """
        slide = self._new_slide()

        # Background — dark teal full bleed
        bg = slide.shapes.add_shape(1, 0, 0, W, H)
        _fill_solid(bg, B.TEAL)
        _no_line(bg)

        # Bottom accent strip
        strip = slide.shapes.add_shape(1, 0, H - Inches(0.5), W, Inches(0.5))
        _fill_solid(strip, B.DARK_TEAL)
        _no_line(strip)

        # Title
        _add_textbox(slide,
                     Inches(0.8), Inches(2.2), W - Inches(1.6), Inches(1.5),
                     text=title, bold=True, size=40, color=B.WHITE,
                     font=B.FONT_HEADING, align='left')

        # Subtitle
        if subtitle:
            _add_textbox(slide,
                         Inches(0.8), Inches(3.9), W - Inches(1.6), Inches(0.8),
                         text=subtitle, size=20, color=B.MINT,
                         font=B.FONT_BODY, align='left')

        # Label (e.g. "Discovery Stage") — top-right badge
        if label:
            _add_textbox(slide,
                         W - Inches(3.5), Inches(0.3), Inches(3.0), Inches(0.4),
                         text=label, size=11, color=B.CYAN,
                         font=B.FONT_BODY, align='right')

        # Product name bottom-left
        _add_textbox(slide,
                     Inches(0.8), H - Inches(0.45), Inches(4), Inches(0.35),
                     text=B.PRODUCT_NAME, size=11, color=B.MINT,
                     font=B.FONT_BODY, align='left')
        return self

    def section_divider(self, title: str, description: str = '') -> 'CB_Pptx':
        """Dark section divider slide."""
        slide = self._new_slide()

        bg = slide.shapes.add_shape(1, 0, 0, W, H)
        _fill_solid(bg, B.DARK_TEAL)
        _no_line(bg)

        # Cyan accent line
        line = slide.shapes.add_shape(1, Inches(0.7), Inches(3.0),
                                       Inches(2), Inches(0.06))
        _fill_solid(line, B.CYAN)
        _no_line(line)

        _add_textbox(slide,
                     Inches(0.7), Inches(2.5), W - Inches(1.4), Inches(1.2),
                     text=title, bold=True, size=36, color=B.WHITE,
                     font=B.FONT_HEADING, align='left')

        if description:
            _add_textbox(slide,
                         Inches(0.7), Inches(3.4), W - Inches(1.4), Inches(1.0),
                         text=description, size=18, color=B.MINT,
                         font=B.FONT_BODY, align='left')

        self._footer(slide)
        return self

    def content(self, title: str, bullets: list[str] = None,
                body_text: str = None, note: str = None) -> 'CB_Pptx':
        """Standard content slide: title bar + bullets or body text."""
        slide = self._new_slide()
        self._title_bar(slide, title)

        top = CONTENT_TOP
        if bullets:
            txBox = slide.shapes.add_textbox(M_LEFT, top, CONTENT_W, CONTENT_H)
            tf = txBox.text_frame
            tf.word_wrap = True
            for idx, b in enumerate(bullets):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                # Detect sub-bullet (starts with spaces or –)
                level = 1 if b.startswith('  ') else 0
                b_clean = b.lstrip()
                bullet_char = '•  ' if level == 0 else '–  '
                bullet_color = B.CYAN if level == 0 else B.TEAL
                _run(p, bullet_char, bold=False, size=16,
                     color=bullet_color, font=B.FONT_BODY)
                _inline_runs(p, b_clean, size=16, color=B.BLACK,
                              font=B.FONT_BODY)
        elif body_text:
            _add_textbox(slide, M_LEFT, top, CONTENT_W, CONTENT_H,
                         text=body_text, size=16, color=B.BLACK,
                         font=B.FONT_BODY, align='left')

        if note:
            _add_textbox(slide,
                         M_LEFT, H - Inches(0.8), CONTENT_W, Inches(0.4),
                         text=f'Note: {note}', size=9, color=B.GRAY,
                         font=B.FONT_BODY, italic=True, align='left')

        self._footer(slide)
        return self

    def two_column(self, title: str,
                   left: list[str], right: list[str],
                   left_label: str = '', right_label: str = '') -> 'CB_Pptx':
        """Two-column slide: findings on left, implications on right."""
        slide = self._new_slide()
        self._title_bar(slide, title)

        col_w   = (CONTENT_W - Inches(0.2)) / 2
        top     = CONTENT_TOP
        left_x  = M_LEFT
        right_x = M_LEFT + col_w + Inches(0.2)

        # Column labels
        if left_label:
            _add_textbox(slide, left_x, top, col_w, Inches(0.35),
                         text=left_label, bold=True, size=13,
                         color=B.TEAL, font=B.FONT_HEADING)
            top += Inches(0.4)

        if right_label:
            _add_textbox(slide, right_x, CONTENT_TOP, col_w, Inches(0.35),
                         text=right_label, bold=True, size=13,
                         color=B.CYAN, font=B.FONT_HEADING)

        # Divider line
        div = slide.shapes.add_shape(
            1, M_LEFT + col_w + Inches(0.07), CONTENT_TOP,
            Inches(0.04), CONTENT_H)
        _fill_solid(div, B.BORDER)
        _no_line(div)

        def _add_bullets(x, items, clr):
            txBox = slide.shapes.add_textbox(x, top, col_w, CONTENT_H)
            tf = txBox.text_frame
            tf.word_wrap = True
            for idx, item in enumerate(items):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                _run(p, '●  ', size=15, color=clr, font=B.FONT_BODY)
                _inline_runs(p, item, size=15, color=B.BLACK, font=B.FONT_BODY)

        _add_bullets(left_x,  left,  B.CYAN)
        _add_bullets(right_x, right, B.TEAL)

        self._footer(slide)
        return self

    def table_slide(self, title: str, headers: list[str],
                    rows: list[list[str]],
                    col_widths: list[float] = None) -> 'CB_Pptx':
        """Slide with a branded data table."""
        from pptx.util import Inches as _I
        slide = self._new_slide()
        self._title_bar(slide, title)

        ncols = len(headers)
        nrows = len(rows) + 1  # +1 for header

        if col_widths:
            cw = [_I(w) for w in col_widths]
        else:
            per = float(CONTENT_W) / ncols
            cw  = [int(per)] * ncols

        top = CONTENT_TOP + _I(0.1)
        row_h = min(_I(0.42), float(CONTENT_H) / nrows)

        tbl_shape = slide.shapes.add_table(
            nrows, ncols,
            M_LEFT, top, CONTENT_W, int(row_h * nrows))
        tbl = tbl_shape.table

        # Column widths
        for i, w in enumerate(cw):
            tbl.columns[i].width = w

        # Header row
        for c, h in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(B.TEAL)
            p = cell.text_frame.paragraphs[0]
            _run(p, h, bold=True, size=12, color=B.WHITE, font=B.FONT_HEADING)

        # Data rows
        for r_idx, row in enumerate(rows):
            bg = B.WHITE if r_idx % 2 == 0 else B.MINT
            for c_idx, val in enumerate(row[:ncols]):
                cell = tbl.cell(r_idx + 1, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(bg)
                p = cell.text_frame.paragraphs[0]
                _inline_runs(p, str(val), size=11, color=B.BLACK,
                              font=B.FONT_BODY)

        self._footer(slide)
        return self

    def quote_slide(self, quote: str, attribution: str = '',
                    context: str = '') -> 'CB_Pptx':
        """Full-bleed quote slide — for research findings / user quotes."""
        slide = self._new_slide()

        bg = slide.shapes.add_shape(1, 0, 0, W, H)
        _fill_solid(bg, B.MINT)
        _no_line(bg)

        # Cyan accent bar left
        bar = slide.shapes.add_shape(1, 0, 0, Inches(0.15), H)
        _fill_solid(bar, B.CYAN)
        _no_line(bar)

        _add_textbox(slide,
                     Inches(0.6), Inches(1.5), W - Inches(1.2), Inches(3.5),
                     text=f'"{quote}"', size=24, color=B.DARK_TEAL,
                     font=B.FONT_BODY, italic=True, align='left')

        if attribution:
            _add_textbox(slide,
                         Inches(0.6), Inches(5.5), W - Inches(1.2), Inches(0.5),
                         text=f'— {attribution}', bold=True, size=13,
                         color=B.TEAL, font=B.FONT_HEADING, align='left')

        if context:
            _add_textbox(slide,
                         Inches(0.6), Inches(6.2), W - Inches(1.2), Inches(0.5),
                         text=context, size=11, color=B.GRAY,
                         font=B.FONT_BODY, align='left')
        return self

    def closing(self, title: str = 'Next Steps',
                bullets: list[str] = None) -> 'CB_Pptx':
        """Closing / CTA slide."""
        slide = self._new_slide()

        bg = slide.shapes.add_shape(1, 0, 0, W, H)
        _fill_solid(bg, B.DARK_TEAL)
        _no_line(bg)

        _add_textbox(slide,
                     Inches(0.8), Inches(1.2), W - Inches(1.6), Inches(1.2),
                     text=title, bold=True, size=36, color=B.WHITE,
                     font=B.FONT_HEADING, align='left')

        if bullets:
            txBox = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.8), W - Inches(1.6), Inches(3.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for idx, b in enumerate(bullets):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                _run(p, '→  ', bold=True, size=18, color=B.CYAN,
                     font=B.FONT_BODY)
                _inline_runs(p, b, size=18, color=B.WHITE, font=B.FONT_BODY)

        _add_textbox(slide,
                     Inches(0.8), H - Inches(0.55), Inches(4), Inches(0.4),
                     text=B.PRODUCT_NAME, size=11, color=B.MINT,
                     font=B.FONT_BODY, align='left')
        return self

    # ── Save ──────────────────────────────────────────────────────────────────

    def save(self, path: str):
        self._prs.save(path)
        n = len(self._prs.slides)
        print(f'Saved: {path}  ({n} slide{"s" if n != 1 else ""})')
        return path

    # ── Markdown builder ──────────────────────────────────────────────────────

    @classmethod
    def build_from_markdown(cls, md_path: str, out_path: str,
                             title: str = None,
                             deck_type: str = 'Summary') -> str:
        """
        Parse a markdown file and generate a best-effort branded slide deck.

        Mapping rules:
          # H1          → cover slide title
          ## H2         → section divider
          ### H3        → content slide title (bullets that follow go on that slide)
          - bullet      → added to current content slide
          > blockquote  → quote slide
          | table |     → table slide (uses the preceding ### as title)
          ---           → ignored
        """
        with open(md_path, 'r', encoding='utf-8') as f:
            content = f.read()

        if not title:
            for line in content.split('\n'):
                s = line.strip()
                if s.startswith('# ') and not s.startswith('## '):
                    title = s[2:].strip()
                    break
            title = title or os.path.splitext(os.path.basename(md_path))[0]

        prs = cls()
        prs.cover(title, subtitle=deck_type)

        lines = content.split('\n')
        i = 0
        current_h3   = None
        current_bullets: list[str] = []
        meta_done    = False

        def _flush():
            nonlocal current_h3, current_bullets
            if current_h3 and current_bullets:
                prs.content(current_h3, bullets=current_bullets)
            elif current_h3:
                prs.content(current_h3, body_text='')
            current_h3      = None
            current_bullets = []

        while i < len(lines):
            line     = lines[i]
            stripped = line.strip()

            if not stripped or stripped == '---':
                i += 1
                continue

            # Skip top metadata block
            if not meta_done:
                if stripped.startswith('# ') and not stripped.startswith('## '):
                    meta_done = True
                    i += 1
                    continue
                if re.match(r'^\*\*[^*]+\*\*', stripped):
                    i += 1
                    continue

            # H2 → section divider
            if stripped.startswith('## ') and not stripped.startswith('### '):
                meta_done = True
                _flush()
                prs.section_divider(stripped[3:].strip())
                i += 1
                continue

            # H3 → new content slide
            if stripped.startswith('### '):
                meta_done = True
                _flush()
                current_h3 = stripped[4:].strip()
                i += 1
                continue

            # Blockquote → quote slide
            if stripped.startswith('>'):
                meta_done = True
                _flush()
                parts = []
                while i < len(lines) and lines[i].strip().startswith('>'):
                    p = lines[i].strip()
                    parts.append('' if p == '>' else p[2:])
                    i += 1
                prs.quote_slide(' '.join(x for x in parts if x))
                continue

            # Table → table slide
            if stripped.startswith('|'):
                meta_done = True
                rows_raw = []
                while i < len(lines) and lines[i].strip().startswith('|'):
                    raw = lines[i].strip()
                    if not re.match(r'^\|[\s\-\|:]+\|$', raw):
                        cells = [c.strip() for c in raw.split('|') if c.strip()]
                        if cells:
                            rows_raw.append(cells)
                    i += 1
                if rows_raw:
                    max_c = max(len(r) for r in rows_raw)
                    rows_raw = [r + [''] * (max_c - len(r)) for r in rows_raw]
                    tbl_title = current_h3 or 'Table'
                    _flush()
                    prs.table_slide(tbl_title, rows_raw[0], rows_raw[1:])
                continue

            # Bullets → accumulate onto current slide
            if stripped.startswith('- ') or stripped.startswith('- [ ]'):
                meta_done = True
                text = stripped[5:].strip() if stripped.startswith('- [ ]') else stripped[2:]
                current_bullets.append(text)
                i += 1
                continue

            # Level-2 indented bullet
            if re.match(r'^[ \t]{3,}- ', line):
                meta_done = True
                current_bullets.append('  ' + re.sub(r'^[ \t]+-\s+', '', line).strip())
                i += 1
                continue

            i += 1

        _flush()
        prs.save(out_path)
        return out_path


# ── CLI ───────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    import argparse
    ap = argparse.ArgumentParser(description='Convert markdown to branded .pptx')
    ap.add_argument('input',          help='Input .md file')
    ap.add_argument('output',         help='Output .pptx file')
    ap.add_argument('--title', '-t',  default=None,      help='Deck title')
    ap.add_argument('--type',  '-y',  default='Summary', help='Deck type label')
    args = ap.parse_args()
    CB_Pptx.build_from_markdown(args.input, args.output,
                                 title=args.title, deck_type=args.type)
